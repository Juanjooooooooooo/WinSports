"""
Generador de la presentación .pptx de WinSports.

Identidad de marca Win Sports: rojo (#E30613), negro y blanco.
Estructura: portada → LEGAL DISCLAIMER → agenda → 21 secciones → cierre.

Uso:
    uv run --with python-pptx python scripts/generar_presentacion.py

Salida:
    docs/WinSports_Presentacion.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "docs" / "WinSports_Presentacion.pptx"

# ── Paleta Win Sports ────────────────────────────────────────────────────────
BLACK = RGBColor(0x0A, 0x0A, 0x0A)        # fondo principal
SURFACE = RGBColor(0x16, 0x16, 0x16)      # tarjetas
SURFACE_2 = RGBColor(0x1F, 0x1F, 0x1F)    # franjas/cabeceras de tarjeta
CODE_BG = RGBColor(0x05, 0x05, 0x05)      # cajas de código
BORDER = RGBColor(0x2D, 0x2D, 0x2D)
BORDER_2 = RGBColor(0x42, 0x42, 0x42)

RED = RGBColor(0xE3, 0x06, 0x13)          # rojo Win Sports
RED_DEEP = RGBColor(0xA0, 0x04, 0x10)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0xFA, 0xFA, 0xFA)
TEXT_DIM = RGBColor(0x9E, 0x9E, 0x9E)
CODE_FG = RGBColor(0xE8, 0xE8, 0xE8)

# Tonos auxiliares (uso muy puntual para diferenciar tarjetas)
SUCCESS = RGBColor(0x4C, 0xAF, 0x82)
WARNING = RGBColor(0xD4, 0xAC, 0x0D)
AMBER = RGBColor(0xE6, 0x7E, 0x22)
INFO = RGBColor(0x4A, 0x9E, 0xFF)
ACCENT = RED                              # alias semántico

SLIDE_W = Cm(33.867)  # 16:9
SLIDE_H = Cm(19.05)

FOOTER_TEXT = "WinSports  ·  Analítica de streaming deportivo"


# ─── Helpers ─────────────────────────────────────────────────────────────────


def make_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    return slide


def paint_background(slide, color=BLACK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(
    slide, text, x, y, w, h, *,
    size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
    italic=False, anchor=MSO_ANCHOR.TOP, font_name="Calibri",
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Cm(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font_name
    return tb


def add_rect(slide, x, y, w, h, fill, line=None, line_width=0.75):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    return rect


def add_pill(slide, text, x, y, w, h, fill, text_color=TEXT, size=11):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    rect.adjustments[0] = 0.5
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.fill.background()
    rect.shadow.inherit = False
    tf = rect.text_frame
    tf.margin_left = tf.margin_right = Cm(0.2)
    tf.margin_top = tf.margin_bottom = Cm(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    r.font.name = "Calibri"


def add_card(slide, title, body_lines, x, y, w, h, *,
             border=BORDER, title_color=RED, body_size=11):
    add_rect(slide, x, y, w, h, SURFACE, border)
    add_rect(slide, x, y, Cm(0.18), h, title_color)

    inner_x = x + Cm(0.55)
    add_text(slide, title, inner_x, y + Cm(0.3), w - Cm(0.7), Cm(1.0),
             size=15, bold=True, color=title_color)
    body = "\n".join(f"  ·  {line}" for line in body_lines)
    add_text(slide, body, inner_x, y + Cm(1.35), w - Cm(0.7), h - Cm(1.5),
             size=body_size, color=TEXT)


def add_section_header(slide, big, subtitle="", number=None):
    add_rect(slide, Cm(0), Cm(0.6), SLIDE_W, Cm(0.08), RED)

    if number is not None:
        add_text(slide, f"{number:02d}", Cm(1.2), Cm(1.0), Cm(3.0), Cm(1.6),
                 size=44, bold=True, color=RED, italic=True)
        add_text(slide, big, Cm(4.0), Cm(1.1), SLIDE_W - Cm(5), Cm(1.6),
                 size=30, bold=True, color=TEXT)
    else:
        add_text(slide, big, Cm(1.2), Cm(1.0), SLIDE_W - Cm(2.4), Cm(1.6),
                 size=32, bold=True, color=TEXT)

    if subtitle:
        add_text(slide, subtitle,
                 Cm(1.2) if number is None else Cm(4.0), Cm(3.0),
                 SLIDE_W - Cm(2.4), Cm(1.0),
                 size=14, color=TEXT_DIM, italic=True)


def add_footer(slide, page_num, total):
    # Marca de la diapositiva: puntito rojo + texto del proyecto
    add_rect(slide, Cm(1.0), SLIDE_H - Cm(0.83), Cm(0.25), Cm(0.25), RED)
    add_text(slide, FOOTER_TEXT, Cm(1.4), SLIDE_H - Cm(0.9), Cm(22), Cm(0.6),
             size=10, color=TEXT_DIM)
    add_text(slide, f"{page_num:02d} / {total:02d}",
             SLIDE_W - Cm(3.5), SLIDE_H - Cm(0.9), Cm(2.7), Cm(0.6),
             size=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def add_code_box(slide, code, x, y, w, h, *, size=10, title=None):
    add_rect(slide, x, y, w, h, CODE_BG, BORDER)
    title_offset = Cm(0)
    if title:
        add_rect(slide, x, y, w, Cm(0.7), SURFACE_2)
        add_text(slide, title, x + Cm(0.4), y + Cm(0.12), w - Cm(0.8), Cm(0.5),
                 size=10, bold=True, color=RED, font_name="Consolas")
        title_offset = Cm(0.7)

    tb = slide.shapes.add_textbox(
        x + Cm(0.3), y + Cm(0.2) + title_offset,
        w - Cm(0.6), h - Cm(0.4) - title_offset,
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Cm(0)
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.name = "Consolas"
        stripped = line.lstrip()
        if stripped.startswith("#") or stripped.startswith("//"):
            r.font.color.rgb = TEXT_DIM
            r.font.italic = True
        else:
            r.font.color.rgb = CODE_FG


def add_arrow(slide, x1, y1, x2, y2, color=RED, width=1.5):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


# ─── Slides ──────────────────────────────────────────────────────────────────


def slide_cover(prs, total):
    slide = add_blank_slide(prs)
    # Franja roja izquierda al estilo broadcast
    add_rect(slide, Cm(0), Cm(0), Cm(2.2), SLIDE_H, RED)
    add_rect(slide, Cm(2.2), Cm(0), Cm(0.15), SLIDE_H, RED_DEEP)

    # Etiqueta superior
    add_text(slide, "INGENIERÍA DE DATOS  ·  PROYECTO FINAL",
             Cm(3.5), Cm(2.3), Cm(20), Cm(0.6),
             size=12, bold=True, color=RED)

    # Identidad grande: WIN (rojo) + SPORTS (blanco)
    add_text(slide, "WIN", Cm(3.3), Cm(3.5), Cm(10), Cm(4.5),
             size=130, bold=True, color=RED)
    add_text(slide, "SPORTS", Cm(13.5), Cm(3.5), Cm(20), Cm(4.5),
             size=130, bold=True, color=WHITE)

    # Tagline
    add_text(slide, "Analítica de streaming deportivo en tiempo real",
             Cm(3.5), Cm(9.5), Cm(28), Cm(1.2),
             size=22, color=RED, italic=True)
    add_text(slide,
             "De un log plano de 300.000 eventos a un dashboard de métricas vivo, "
             "sobre MongoDB Atlas, FastAPI y React",
             Cm(3.5), Cm(11.0), Cm(28), Cm(1.5),
             size=14, color=TEXT_DIM)

    # Badges del stack
    badges = [
        ("MongoDB Atlas", RED),
        ("Motor async", RED_DEEP),
        ("FastAPI", SURFACE_2),
        ("Pydantic v2", SURFACE_2),
        ("React 19 + Vite", SURFACE_2),
        ("Recharts", SURFACE_2),
        ("Leaflet", SURFACE_2),
        ("uv · Python 3.14", SURFACE_2),
    ]
    x_start = Cm(3.5)
    y = Cm(13.6)
    for label, color in badges:
        add_pill(slide, label, x_start, y, Cm(3.4), Cm(0.95), color, WHITE, 11)
        x_start += Cm(3.6)

    # Pie de portada
    add_rect(slide, Cm(3.5), SLIDE_H - Cm(2.5), Cm(28), Cm(0.06), RED)
    add_text(slide,
             f"{total} diapositivas  ·  WinSports v0.1.0  ·  Identidad institucional Win Sports",
             Cm(3.5), SLIDE_H - Cm(2.0), Cm(28), Cm(1),
             size=12, color=TEXT_DIM, italic=True)


def slide_legal_disclaimer(prs, total, page):
    slide = add_blank_slide(prs)
    # Marco rojo grueso característico de aviso legal
    add_rect(slide, Cm(0), Cm(0), SLIDE_W, Cm(0.6), RED)
    add_rect(slide, Cm(0), SLIDE_H - Cm(0.6), SLIDE_W, Cm(0.6), RED)
    add_rect(slide, Cm(0), Cm(0), Cm(0.6), SLIDE_H, RED)
    add_rect(slide, SLIDE_W - Cm(0.6), Cm(0), Cm(0.6), SLIDE_H, RED)

    # Etiqueta superior
    add_text(slide, "AVISO  ·  LECTURA OBLIGATORIA ANTES DE CONTINUAR",
             Cm(1.5), Cm(1.0), SLIDE_W - Cm(3), Cm(0.7),
             size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # Título principal
    add_text(slide, "LEGAL DISCLAIMER",
             Cm(1.5), Cm(2.0), SLIDE_W - Cm(3), Cm(2.5),
             size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Mensaje central — el texto literal que el ponente lee en voz alta
    add_rect(slide, Cm(2.5), Cm(5.5), SLIDE_W - Cm(5), Cm(4.5), SURFACE, RED, 1.5)
    add_text(slide,
             "En esta presentación se mostrará información sensible "
             "protegida bajo secreto industrial.",
             Cm(3.5), Cm(6.0), SLIDE_W - Cm(7), Cm(2.0),
             size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide,
             "Se delega a discreción por parte de los asistentes "
             "el resguardo y manejo responsable del contenido aquí expuesto.",
             Cm(3.5), Cm(8.0), SLIDE_W - Cm(7), Cm(1.8),
             size=18, italic=True, color=RED, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # Tres tarjetas inferiores con las condiciones específicas
    cards = [
        (
            "Confidencialidad",
            [
                "No grabar, no fotografiar, no transmitir pantallas.",
                "No redistribuir el material fuera del recinto.",
                "El uso es exclusivo para esta sesión.",
            ],
            RED,
        ),
        (
            "Naturaleza académica",
            [
                "Proyecto con fines educativos del curso.",
                "No constituye producto ni servicio en producción.",
                "Métricas para ilustrar; no para decidir negocio.",
            ],
            RED_DEEP,
        ),
        (
            "Marcas y datos",
            [
                "Win Sports y los nombres de eventos/ligas son de",
                "  sus respectivos titulares; sin afiliación.",
                "Identificadores anonimizados (prefijos), geo aproximada.",
            ],
            WHITE,
        ),
    ]
    for i, (title, items, color) in enumerate(cards):
        x = Cm(2.5) + i * Cm(9.6)
        y = Cm(10.5)
        add_rect(slide, x, y, Cm(9), Cm(6.5), SURFACE, BORDER_2)
        add_rect(slide, x, y, Cm(9), Cm(0.6), color)
        text_color = WHITE if color == WHITE else WHITE
        add_text(slide, title, x, y + Cm(0.1), Cm(9), Cm(0.6),
                 size=13, bold=True, color=BLACK if color == WHITE else WHITE,
                 align=PP_ALIGN.CENTER)
        for j, b in enumerate(items):
            add_text(slide, "·  " + b,
                     x + Cm(0.4), y + Cm(1.1) + j * Cm(0.95),
                     Cm(8.3), Cm(1.0),
                     size=11, color=TEXT)

    add_footer(slide, page, total)


def slide_agenda(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Agenda", "Veintiún paradas por el sistema completo", number=0)

    items = [
        ("01", "El problema y el dataset", RED),
        ("02", "Stack técnico", RED),
        ("03", "Arquitectura de datos", RED),
        ("04", "Modelo NoSQL  -  3 colecciones", RED),
        ("05", "Colección events", RED),
        ("06", "Colección sessions", RED),
        ("07", "Colección content_stats", RED),
        ("08", "Validadores $jsonSchema", RED),
        ("09", "Índices (16 en total)", RED),
        ("10", "Ingesta CSV → events", RED),
        ("11", "Pipeline: reconstrucción de sesiones", RED),
        ("12", "Pipeline: content_stats", RED),
        ("13", "Alertas en vivo (rojo · amarillo · azul)", RED),
        ("14", "Operadores de agregación", RED),
        ("15", "Change Streams  -  triggers", RED),
        ("16", "API REST  -  18 endpoints", RED),
        ("17", "Dashboard · Overview", RED),
        ("18", "Dashboard · QoE", RED),
        ("19", "Dashboard · Usuarios y Alertas", RED),
        ("20", "Dashboard · Admin (CSV + edición)", RED),
        ("21", "Tests, calidad y números", RED),
    ]

    cols = 3
    col_w = Cm(10.0)
    rows = (len(items) + cols - 1) // cols
    for i, (n, label, color) in enumerate(items):
        col = i % cols
        row = i // cols
        x = Cm(1.5) + col * (col_w + Cm(0.5))
        y = Cm(4.3) + row * Cm(1.55)
        add_rect(slide, x, y, col_w, Cm(1.25), SURFACE, BORDER)
        add_rect(slide, x, y, Cm(0.15), Cm(1.25), color)
        add_text(slide, n, x + Cm(0.5), y + Cm(0.25), Cm(1.3), Cm(0.8),
                 size=14, bold=True, color=color, font_name="Consolas")
        add_text(slide, label, x + Cm(1.9), y + Cm(0.27),
                 col_w - Cm(2.2), Cm(0.85),
                 size=12, color=TEXT)

    add_footer(slide, page, total)


def slide_problema(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "El problema y el dataset",
                       "Un log plano de eventos de reproducción no responde por sí solo",
                       number=1)

    add_card(slide, "El dataset",
             [
                 "CSV con 22 columnas por fila (Date, CountryCode, SubscriberID…)",
                 "Eventos crudos: START, PAUSE, RESUME, RE-BUFFERING, COMPLETE…",
                 "Cobertura geográfica: 13 países (CO mayoritario)",
                 "Tres tipos de dispositivo permitidos: ANDR, IOS, WEB",
                 "Quirk conocido: Position y Duration en unidades mixtas",
                 "Tamaño actual cargado: ~300.000 eventos",
             ],
             Cm(1.5), Cm(4.5), Cm(15), Cm(6.5),
             title_color=RED, body_size=12)

    add_card(slide, "Preguntas de negocio",
             [
                 "¿Qué contenido se ve más y qué se termina más?",
                 "¿Cuál es la calidad de experiencia (buffer, rebuffering)?",
                 "¿Quiénes son los Heavy / Mid / Vague users?",
                 "¿Dónde y a qué horas se concentra la actividad?",
                 "¿Qué contenidos están degradando el servicio ahora mismo?",
                 "¿Cuáles dispositivos pesan más en el catálogo?",
             ],
             Cm(17.5), Cm(4.5), Cm(15), Cm(6.5),
             title_color=RED_DEEP, body_size=12)

    add_card(slide, "Reto técnico",
             [
                 "Convertir un log inmutable en sesiones de reproducción",
                 "Derivar métricas por contenido sin duplicar datos",
                 "Mantener todo sincronizado al ingerir CSVs nuevos",
                 "Exponerlo en una API limpia consumida por el dashboard",
                 "Mantener la API responsiva incluso con cientos de miles de docs",
             ],
             Cm(1.5), Cm(11.3), Cm(31), Cm(4.5),
             title_color=WHITE, body_size=12)

    add_rect(slide, Cm(1.5), Cm(16.2), SLIDE_W - Cm(3), Cm(1.4), SURFACE, RED, 1.0)
    add_text(slide,
             '"Que MongoDB haga el trabajo pesado: agregar dentro de la base, no fuera."',
             Cm(1.5), Cm(16.3), SLIDE_W - Cm(3), Cm(1.2),
             size=16, bold=True, color=RED, align=PP_ALIGN.CENTER,
             italic=True, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, page, total)


def slide_stack(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Stack técnico",
                       "Tres capas modulares: la API jamás toca Mongo directo; el front solo consume REST",
                       number=2)

    cols = [
        ("BACKEND", RED, [
            ("FastAPI + Uvicorn", "API REST async, OpenAPI/Swagger autogenerado"),
            ("Motor 3.x", "Driver async oficial de MongoDB"),
            ("Pydantic v2", "Schemas tipados de request/response"),
            ("pydantic-settings", "Configuración tipada desde .env"),
            ("PyMongo bulk_write", "Upserts idempotentes (UpdateOne)"),
            ("pandas + math", "Parseo y casteo de filas del CSV"),
            ("pytest + mongomock", "81 tests del contrato y los pipelines"),
            ("Python 3.14 + uv", "Entorno reproducible, lock file"),
        ]),
        ("BASE DE DATOS", RED_DEEP, [
            ("MongoDB Atlas", "Replica set: habilita Change Streams"),
            ("$jsonSchema", 'validationLevel="moderate", action="warn"'),
            ("Aggregation Framework", "$group, $switch, $hour, $dayOfWeek, $cond"),
            ("Change Streams", "Sincroniza derivadas en tiempo real"),
            ("Índices compuestos", "Reconstrucción y rankings rápidos"),
            ("Índice único", "Clave anti-duplicados en sessions"),
            ("collMod", "Re-aplica validators de forma idempotente"),
            ("3 colecciones · 16 índices", "Modelo afinado al uso real"),
        ]),
        ("FRONTEND", SURFACE_2, [
            ("React 19 + Vite", "SPA con proxy /api → :8000"),
            ("Recharts", "Barras, líneas, funnel y heatmap"),
            ("Leaflet + react-leaflet", "Mapa de burbujas de usuarios"),
            ("Vitest + Testing Library", "23 tests de componentes"),
            ("CSS theme custom", "Tema oscuro unificado"),
            ("Fetch API nativa", "Consumo de la REST sin SDK"),
            ("eslint + react-refresh", "Hot reload + linting"),
            ("Node 18+", "Toolchain moderno"),
        ]),
    ]

    for i, (title, header_color, items) in enumerate(cols):
        x = Cm(1.5) + i * Cm(10.6)
        add_rect(slide, x, Cm(4.5), Cm(10), Cm(13.5), SURFACE, BORDER)
        add_rect(slide, x, Cm(4.5), Cm(10), Cm(1.0), header_color)
        title_color = BLACK if header_color == SURFACE_2 else WHITE
        if header_color == SURFACE_2:
            title_color = WHITE
        add_text(slide, title, x, Cm(4.65), Cm(10), Cm(0.8),
                 size=14, bold=True, color=title_color, align=PP_ALIGN.CENTER)
        for j, (tech, desc) in enumerate(items):
            yy = Cm(6.0) + j * Cm(1.45)
            add_text(slide, tech, x + Cm(0.5), yy, Cm(9), Cm(0.55),
                     size=12, bold=True, color=WHITE)
            add_text(slide, desc, x + Cm(0.5), yy + Cm(0.55), Cm(9), Cm(0.7),
                     size=10, color=TEXT_DIM, italic=True)

    add_footer(slide, page, total)


def slide_arquitectura(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Arquitectura de datos",
                       "Flujo unidireccional: el CSV es la fuente, todo lo demás se deriva",
                       number=3)

    box_w, box_h = Cm(7), Cm(3.4)
    y_top = Cm(5.0)
    xs = [Cm(1.5), Cm(10.0), Cm(18.5), Cm(27.0) - Cm(1.5)]

    boxes = [
        ("CSV", AMBER, "data/*.csv\nEventos crudos\nde reproducción"),
        ("events", RED, "Fuente de verdad\n1 doc por evento\n~300k docs"),
        ("sessions", RED_DEEP, "Derivada\nAgrupa por usuario,\nparte por START"),
        ("content_stats", WHITE, "Derivada\nMétricas agregadas\npor contenido"),
    ]
    for (name, color, desc), x in zip(boxes, xs):
        add_rect(slide, x, y_top, box_w, box_h, SURFACE, color)
        title_color = color if color != WHITE else WHITE
        add_text(slide, name, x, y_top + Cm(0.3), box_w, Cm(1),
                 size=17, bold=True, color=title_color, align=PP_ALIGN.CENTER,
                 font_name="Consolas")
        add_text(slide, desc, x + Cm(0.3), y_top + Cm(1.3), box_w - Cm(0.6), Cm(2),
                 size=11, color=TEXT, align=PP_ALIGN.CENTER)

    labels = ["load_csv.py / upload-csv", "change stream insert", "change stream insert+update"]
    for i in range(3):
        x1 = xs[i] + box_w
        x2 = xs[i + 1]
        add_arrow(slide, x1, y_top + Cm(1.7), x2, y_top + Cm(1.7), RED)
        add_text(slide, labels[i], x1, y_top - Cm(0.1), x2 - x1, Cm(0.6),
                 size=9, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER)

    # Bloque "qué corre cuándo"
    add_rect(slide, Cm(1.5), Cm(9.5), SLIDE_W - Cm(3), Cm(3.2), SURFACE, BORDER_2)
    add_text(slide, "Núcleo orientado a eventos", Cm(2.0), Cm(9.7),
             Cm(20), Cm(0.7), size=14, bold=True, color=RED)
    add_text(slide,
             "En PROD, los Change Streams de MongoDB reaccionan a cada insert/update y "
             "reconstruyen 'sessions' y 'content_stats' sin polling.\n"
             "En DEV, las derivadas se construyen en lote con scripts/test_pipeline.py o "
             "desde el botón 'Subir CSV' en la pestaña Admin (que también re-corre los pipelines).\n"
             "_safe_watch() en api/main.py degrada con gracia si el cluster no soporta change streams: "
             "loguea el motivo y la API sigue viva sirviendo las derivadas existentes.",
             Cm(2.0), Cm(10.4), SLIDE_W - Cm(4), Cm(2.4),
             size=11, color=TEXT)

    add_card(slide, "Por qué derivar (y no recalcular en cada request)",
             [
                 "El log crudo es inmutable: auditable y rebobinable",
                 "Las métricas viven pre-agregadas → el dashboard responde rápido",
                 "Re-ingerir un CSV actualiza la sesión vía upsert, no la duplica",
                 "Una caída del Change Stream NO degrada las páginas (lo derivado sigue)",
             ],
             Cm(1.5), Cm(13.0), SLIDE_W - Cm(3), Cm(4.5),
             title_color=RED, body_size=12)

    add_footer(slide, page, total)


def slide_modelo_overview(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Modelo NoSQL  -  tres colecciones",
                       "Una fuente de verdad, dos derivadas. Sin joins en runtime.",
                       number=4)

    cards = [
        ("events",
         ["Log crudo del CSV  ·  ~300.000 docs",
          "Un evento de reproducción por documento",
          "Inmutable: nunca se modifica, solo se inserta",
          "Validator $jsonSchema flexible (warn)",
          "Índices de reconstrucción + filtros temporales"],
         RED),
        ("sessions",
         ["Sesiones de reproducción  ·  ~28.500 docs",
          "Una sesión = todos los eventos entre un START y el siguiente",
          "Métricas QoE pre-calculadas (buffer, rebuffer, completion)",
          "Hitos del funnel (firstquartile · midpoint · thirdquartile · complete)",
          "Índice ÚNICO: (subscriber_id, customer_id, start_time)"],
         RED_DEEP),
        ("content_stats",
         ["Métricas por contenido  ·  ~10.800 docs",
          "Una pieza puntual = un customer_id",
          "total_plays, unique_viewers, completion_rate, avg_buffer_time",
          "Tasas del funnel (firstquartile_rate … completion_rate)",
          "Breakdown por dispositivo  ·  last_updated UTC"],
         WHITE),
    ]
    for i, (name, items, color) in enumerate(cards):
        x = Cm(1.5) + i * Cm(10.6)
        add_rect(slide, x, Cm(4.5), Cm(10), Cm(11.5), SURFACE, BORDER)
        add_rect(slide, x, Cm(4.5), Cm(10), Cm(0.18), color)
        add_text(slide, name, x + Cm(0.5), Cm(4.9), Cm(9), Cm(0.9),
                 size=22, bold=True, color=color, font_name="Consolas")
        for j, it in enumerate(items):
            add_text(slide, "·  " + it, x + Cm(0.5), Cm(6.4) + j * Cm(1.65),
                     Cm(9), Cm(1.5),
                     size=12, color=TEXT)

    # Diagrama lateral del flujo
    add_rect(slide, Cm(1.5), Cm(16.4), SLIDE_W - Cm(3), Cm(1.2), SURFACE, RED, 1.0)
    add_text(slide,
             "CSV  →  load_csv.py  →  events  ──(change stream)→  sessions  ──(change stream)→  content_stats",
             Cm(1.5), Cm(16.5), SLIDE_W - Cm(3), Cm(1.0),
             size=14, bold=True, color=RED, align=PP_ALIGN.CENTER, italic=True,
             anchor=MSO_ANCHOR.MIDDLE, font_name="Consolas")

    add_footer(slide, page, total)


def slide_coleccion_events(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Colección events  -  fuente de verdad",
                       "Un documento por cada evento de reproducción del CSV",
                       number=5)

    code = (
        "EVENTS_VALIDATOR = {\n"
        '    "$jsonSchema": {\n'
        '        "bsonType": "object",\n'
        '        "required": ["date", "subscriber_id",\n'
        '                     "customer_id", "type_event",\n'
        '                     "device_type"],\n'
        '        "properties": {\n'
        '            "date":          {"bsonType": "date"},\n'
        '            "subscriber_id": {"bsonType": "string"},\n'
        '            "customer_id":   {"bsonType": "string"},\n'
        '            "type_event":    {"bsonType": "string"},\n'
        '            "device_type":   {"bsonType": "string",\n'
        '                              "enum": ["ANDR","IOS","WEB"]},\n'
        '            "country_code":  {"bsonType": ["string","null"]},\n'
        '            "title":         {"bsonType": ["string","null"]},\n'
        '            "series_title":  {"bsonType": ["string","null"]},\n'
        '            "content_type":  {"bsonType": ["string","null"]},\n'
        '            "genres":        {"bsonType": ["string","null"]},\n'
        '            # bitrate/position superan int32 → int + long\n'
        '            "bitrate":       {"bsonType": ["int","long","null"]},\n'
        '            "position":      {"bsonType": ["int","long","null"]},\n'
        '            "buffer_time":   {"bsonType": ["double","int","null"]},\n'
        '            "duration":      {"bsonType": ["int","long","null"]},\n'
        "        },\n"
        "    }\n"
        "}"
    )
    add_code_box(slide, code, Cm(1.5), Cm(4.5), Cm(19.5), Cm(13),
                 size=10, title="db/collections/events.py")

    add_card(slide, "Decisiones de diseño",
             [
                 "Sin additionalProperties:False → esquema flexible",
                 'validationAction="warn": no para el batch de 300k filas',
                 "Enteros admiten int Y long (bitrate ≫ 2³¹ posible)",
                 "country_code y title aceptan null (algunos eventos no lo traen)",
                 "buffer_time como double + int + null (mix del dataset)",
                 "Re-aplicable con collMod (idempotente al re-arrancar)",
             ],
             Cm(22.0), Cm(4.5), Cm(10.5), Cm(8),
             title_color=RED, body_size=11)

    add_card(slide, "Índices",
             [
                 "idx_events_session_rebuild (subscriber, customer, date)",
                 "idx_events_date (date desc)",
                 "idx_events_type_date (type_event, date desc)",
                 "idx_events_content_type (customer_id, type_event)",
                 "idx_events_device (device_type)",
             ],
             Cm(22.0), Cm(12.8), Cm(10.5), Cm(4.7),
             title_color=RED_DEEP, body_size=11)

    add_footer(slide, page, total)


def slide_coleccion_sessions(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Colección sessions  -  sesiones de reproducción",
                       "Una sesión = todos los eventos entre dos START consecutivos",
                       number=6)

    code = (
        "SESSIONS_VALIDATOR = {\n"
        '    "$jsonSchema": {\n'
        '        "bsonType": "object",\n'
        '        "required": ["subscriber_id", "customer_id",\n'
        '                     "start_time", "device_type"],\n'
        '        "properties": {\n'
        '            "subscriber_id":    {"bsonType": "string"},\n'
        '            "customer_id":      {"bsonType": "string"},\n'
        '            "start_time":       {"bsonType": "date"},\n'
        '            "end_time":         {"bsonType": ["date","null"]},\n'
        '            "device_type":      {"bsonType": "string",\n'
        '                                 "enum": ["ANDR","IOS","WEB"]},\n'
        '            "total_buffer_time":{"bsonType": ["double","null"]},\n'
        '            "rebuffer_count":   {"bsonType": ["int","null"]},\n'
        '            "pause_count":      {"bsonType": ["int","null"]},\n'
        '            "max_position":     {"bsonType": ["int","null"]},\n'
        '            "completion_pct":   {"bsonType": ["double","null"]},\n'
        '            "reached_firstquartile":  {"bsonType": "bool"},\n'
        '            "reached_midpoint":       {"bsonType": "bool"},\n'
        '            "reached_thirdquartile":  {"bsonType": "bool"},\n'
        '            "reached_complete":       {"bsonType": "bool"},\n'
        '            "dominant_bitrate_type":  {"bsonType":["string","null"]},\n'
        "        }}}"
    )
    add_code_box(slide, code, Cm(1.5), Cm(4.5), Cm(19.5), Cm(13),
                 size=10, title="db/collections/sessions.py")

    add_card(slide, "Clave de negocio (ÚNICA)",
             [
                 "(subscriber_id, customer_id, start_time)",
                 "Índice unique=True respaldando la clave",
                 "Todos los escritos usan upsert con esa clave",
                 "Re-procesar un CSV ACTUALIZA, no duplica",
                 "Las sesiones tienen identidad estable",
             ],
             Cm(22.0), Cm(4.5), Cm(10.5), Cm(6.0),
             title_color=RED, body_size=11)

    add_card(slide, "Métricas pre-calculadas",
             [
                 "total_buffer_time, rebuffer_count, pause_count",
                 "max_position y completion_pct (capado a 100%)",
                 "Funnel: reached_firstquartile … reached_complete",
                 "dominant_bitrate_type (HIGH/MID/LOW dominante)",
                 "event_count para perfil del usuario",
             ],
             Cm(22.0), Cm(10.8), Cm(10.5), Cm(6.7),
             title_color=RED_DEEP, body_size=11)

    add_footer(slide, page, total)


def slide_coleccion_content_stats(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Colección content_stats  -  métricas por contenido",
                       "Un documento por customer_id; el dashboard la consulta lista",
                       number=7)

    code = (
        "CONTENT_STATS_VALIDATOR = {\n"
        '    "$jsonSchema": {\n'
        '        "bsonType": "object",\n'
        '        "required": ["customer_id", "title", "total_plays"],\n'
        '        "properties": {\n'
        '            "customer_id":   {"bsonType": "string"},\n'
        '            "title":         {"bsonType": "string"},\n'
        '            "series_title":  {"bsonType": ["string","null"]},\n'
        '            "content_type":  {"bsonType": ["string","null"]},\n'
        '            "genres":        {"bsonType": ["string","null"]},\n'
        '            "release_year":  {"bsonType": ["int","null"]},\n'
        '            "duration":      {"bsonType": ["int","null"]},\n'
        '            # Volumen\n'
        '            "total_plays":    {"bsonType": "int"},\n'
        '            "unique_viewers": {"bsonType": ["int","null"]},\n'
        '            # Funnel\n'
        '            "firstquartile_rate":{"bsonType":["double","null"]},\n'
        '            "midpoint_rate":     {"bsonType":["double","null"]},\n'
        '            "thirdquartile_rate":{"bsonType":["double","null"]},\n'
        '            "completion_rate":   {"bsonType":["double","null"]},\n'
        '            # QoE\n'
        '            "avg_buffer_time":   {"bsonType":["double","null"]},\n'
        '            "rebuffer_rate":     {"bsonType":["double","null"]},\n'
        '            # Breakdown {ANDR:N, IOS:N, WEB:N}\n'
        '            "plays_by_device":   {"bsonType":["object","null"]},\n'
        '            "last_updated":      {"bsonType":["date","null"]},\n'
        "        }}}"
    )
    add_code_box(slide, code, Cm(1.5), Cm(4.5), Cm(19.5), Cm(13),
                 size=10, title="db/collections/content_stats.py")

    add_card(slide, "Índices (5 en total)",
             [
                 "idx_content_stats_id  (customer_id, UNIQUE)",
                 "idx_content_stats_plays  (total_plays desc)",
                 "idx_content_stats_buffer (avg_buffer_time desc)",
                 "idx_content_stats_genre  (genres, total_plays desc)",
                 "idx_content_stats_series (series_title, total_plays)",
             ],
             Cm(22.0), Cm(4.5), Cm(10.5), Cm(6.0),
             title_color=RED, body_size=11)

    add_card(slide, "Para qué la usa el dashboard",
             [
                 "Top contenido y rankings de Overview",
                 "Buffer por contenido y QoE (peores primero)",
                 "Completion ranking de Usuarios",
                 "Filtros por género y agrupación por serie",
                 "last_updated permite auditar la frescura",
             ],
             Cm(22.0), Cm(10.8), Cm(10.5), Cm(6.7),
             title_color=RED_DEEP, body_size=11)

    add_footer(slide, page, total)


def slide_validators(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Validadores $jsonSchema",
                       'Nivel "moderate" + acción "warn": valida sin frenar la carga',
                       number=8)

    add_card(slide, "Por qué moderate + warn",
             [
                 'moderate: aplica el validator a inserts y a updates',
                 '  pero PERMITE updates que dejen el doc fuera del schema',
                 '  → útil cuando se evoluciona el modelo sin migrar todo',
                 'warn: una fila rara loguea pero NO frena el batch',
                 '  → al cargar 300k filas, una mala no cancela el proceso',
                 'collMod permite re-aplicar el validator sin recrear la',
                 '  colección — idempotente al arrancar la API',
             ],
             Cm(1.5), Cm(4.5), Cm(15), Cm(7.5),
             title_color=RED, body_size=12)

    add_card(slide, "Decisiones específicas de WinSports",
             [
                 "events SIN additionalProperties:False (esquema flexible)",
                 "Enteros aceptan int + long (bitrate/position grandes)",
                 "Tipos null explícitos en campos opcionales",
                 "device_type con enum cerrado (ANDR/IOS/WEB)",
                 "Pydantic v2 refleja la misma forma en la API",
                 "Si el cluster está vacío, setup_*_collection crea+valida",
             ],
             Cm(17.5), Cm(4.5), Cm(15), Cm(7.5),
             title_color=RED_DEEP, body_size=12)

    code = (
        "await db.create_collection(\n"
        '    "events",\n'
        "    validator=EVENTS_VALIDATOR,\n"
        '    validationLevel="moderate",   # inserts + updates\n'
        '    validationAction="warn",      # loguea, no rechaza\n'
        ")\n\n"
        "# Idempotente: si ya existe, re-aplica el validator\n"
        'await db.command("collMod", "events",\n'
        "                 validator=EVENTS_VALIDATOR,\n"
        '                 validationLevel="moderate",\n'
        '                 validationAction="warn")'
    )
    add_code_box(slide, code, Cm(1.5), Cm(12.5), SLIDE_W - Cm(3), Cm(5),
                 size=11, title="db/collections/events.py · setup_events_collection()")

    add_footer(slide, page, total)


def slide_indices(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Índices  -  16 en total entre las tres colecciones",
                       "Compuestos para reconstrucción y rankings; único para blindar la clave de sesión",
                       number=9)

    code = (
        "# events (5)\n"
        "IndexModel([(\"subscriber_id\", ASCENDING),\n"
        "            (\"customer_id\", ASCENDING),\n"
        "            (\"date\", ASCENDING)], name=\"idx_events_session_rebuild\")\n"
        "IndexModel([(\"date\", DESCENDING)], name=\"idx_events_date\")\n"
        "IndexModel([(\"type_event\", ASCENDING),\n"
        "            (\"date\", DESCENDING)], name=\"idx_events_type_date\")\n"
        "IndexModel([(\"customer_id\", ASCENDING),\n"
        "            (\"type_event\", ASCENDING)], name=\"idx_events_content_type\")\n"
        "IndexModel([(\"device_type\", ASCENDING)], name=\"idx_events_device\")\n\n"
        "# sessions (6)  -  el unique=True es CLAVE del modelo\n"
        "IndexModel([(\"subscriber_id\", ASCENDING),\n"
        "            (\"customer_id\", ASCENDING),\n"
        "            (\"start_time\", ASCENDING)],\n"
        "           name=\"idx_sessions_unique_key\", unique=True)\n"
        "IndexModel([(\"start_time\", DESCENDING),\n"
        "            (\"total_buffer_time\", DESCENDING)],\n"
        "           name=\"idx_sessions_time_buffer\")\n"
        "IndexModel([(\"customer_id\", ASCENDING),\n"
        "            (\"reached_complete\", ASCENDING)],\n"
        "           name=\"idx_sessions_content_complete\")\n\n"
        "# content_stats (5)\n"
        "IndexModel([(\"customer_id\", ASCENDING)],\n"
        "           name=\"idx_content_stats_id\", unique=True)\n"
        "IndexModel([(\"total_plays\", DESCENDING)],\n"
        "           name=\"idx_content_stats_plays\")"
    )
    add_code_box(slide, code, Cm(1.5), Cm(4.5), Cm(20), Cm(13),
                 size=9, title="db/indexes/{events,sessions,content_stats}.py")

    bullets = [
        ("Reconstrucción", "(subscriber, customer, date) trae los eventos en orden cronológico"),
        ("Único", "(subscriber, customer, start_time) → upsert sin duplicar sesiones"),
        ("Temporal", "(start_time desc, buffer desc) alimenta alertas recientes"),
        ("QoE", "(type_event, date desc) filtra RE-BUFFERING / START-BUFFERING"),
        ("Funnel", "(customer_id, reached_complete) cuenta sesiones completas"),
        ("Rankings", "device_type, customer_id, genres para breakdowns rápidos"),
    ]
    for i, (title, desc) in enumerate(bullets):
        y = Cm(4.7) + i * Cm(2.1)
        add_rect(slide, Cm(22.5), y, Cm(10), Cm(1.85), SURFACE, BORDER)
        add_rect(slide, Cm(22.5), y, Cm(0.15), Cm(1.85), RED)
        add_text(slide, title, Cm(22.9), y + Cm(0.15), Cm(9), Cm(0.6),
                 size=13, bold=True, color=RED)
        add_text(slide, desc, Cm(22.9), y + Cm(0.85), Cm(9), Cm(1.0),
                 size=10, color=TEXT)

    add_footer(slide, page, total)


def slide_ingesta_csv(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Ingesta CSV → events",
                       "Parseo, casteo a bsonType, filtrado y batch insert idempotente",
                       number=10)

    add_card(slide, "Pipeline de carga",
             [
                 "pandas.read_csv(BytesIO(raw), dtype=str) — todo string primero",
                 "parse_row(): mapea Date,SubscriberID,… → date,subscriber_id,…",
                 "Casteo a bsonType: _to_date / _to_int / _to_float / _to_str",
                 "NaN de pandas → None (compatible con bsonType: null)",
                 "Filtra filas SIN campos required (REQUIRED_FIELDS)",
                 "insert_many en batches de 1.000 con ordered=False",
                 "BulkWriteError: cuenta los insertados reales con nInserted",
             ],
             Cm(1.5), Cm(4.5), Cm(15.5), Cm(8),
             title_color=RED, body_size=12)

    add_card(slide, "Re-uso en CLI y en Admin",
             [
                 "scripts/load_csv.py (terminal)  → llama a parse_records + insert_events",
                 "POST /api/admin/upload-csv     → llama EXACTAMENTE lo mismo",
                 "ingest_csv() encadena: events → build_sessions → build_content_stats",
                 "El usuario sube un CSV desde el dashboard y todo se rehidrata solo",
                 "Después responde: rows_total, rows_skipped, inserted, sessions_built, …",
             ],
             Cm(17.5), Cm(4.5), Cm(15), Cm(8),
             title_color=RED_DEEP, body_size=12)

    code = (
        "def parse_row(row: dict) -> dict:\n"
        "    return {\n"
        "        \"date\":          _to_date(row.get(\"Date\")),\n"
        "        \"subscriber_id\": _to_str(row.get(\"SubscriberID\")),\n"
        "        \"customer_id\":   _to_str(row.get(\"CustomerId\")),\n"
        "        \"type_event\":    _to_str(row.get(\"TypeEvent\")),\n"
        "        \"device_type\":   _to_str(row.get(\"DeviceType\")),\n"
        "        \"bitrate\":       _to_int(row.get(\"Bitrate\")),\n"
        "        \"position\":      _to_int(row.get(\"Position\")),\n"
        "        \"buffer_time\":   _to_float(row.get(\"BufferTime\")),\n"
        "        \"duration\":      _to_int(row.get(\"Duration\")),\n"
        "        # … 14 campos más\n"
        "    }\n\n"
        "async def insert_events(db, docs):\n"
        "    for i in range(0, len(docs), BATCH_SIZE):\n"
        "        batch = docs[i:i + BATCH_SIZE]\n"
        "        try:\n"
        "            r = await db.events.insert_many(batch, ordered=False)\n"
        "            inserted += len(r.inserted_ids)\n"
        "        except Exception as e:  # BulkWriteError\n"
        "            inserted += e.details.get(\"nInserted\", 0)"
    )
    add_code_box(slide, code, Cm(1.5), Cm(12.8), SLIDE_W - Cm(3), Cm(4.8),
                 size=10, title="db/csv_ingest.py")

    add_footer(slide, page, total)


def slide_pipeline_sessions(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Pipeline: reconstrucción de sesiones",
                       "De un log plano a sesiones con métricas, en una sola pasada por usuario",
                       number=11)

    add_card(slide, "Algoritmo",
             [
                 "1. distinct subscriber_id sobre events",
                 "2. Para cada subscriber: traer eventos ORDENADOS por date",
                 "3. Agrupar por (subscriber_id, customer_id)",
                 "4. Partir en sub-sesiones cada vez que aparece un START",
                 "5. Métricas por sesión: buffer, rebuffer, pause, completion",
                 "6. bulk_write con UpdateOne(..., upsert=True) por clave única",
                 "7. ordered=False: los upserts no se bloquean entre sí",
             ],
             Cm(1.5), Cm(4.5), Cm(13.5), Cm(7),
             title_color=RED, body_size=12)

    add_card(slide, "Garantías",
             [
                 "Sin duplicados (índice único respalda)",
                 "Idempotente: re-procesar actualiza",
                 "completion_pct capado a 100%",
                 "  (unidades mixtas en el dataset)",
                 "Proyección de campos en el find",
                 "  → no se trae el documento entero",
             ],
             Cm(1.5), Cm(11.6), Cm(13.5), Cm(5.7),
             title_color=RED_DEEP, body_size=12)

    code = (
        "def _build_sessions(events):\n"
        "    groups = defaultdict(list)\n"
        "    for ev in events:\n"
        "        groups[(ev['subscriber_id'], ev['customer_id'])].append(ev)\n\n"
        "    for evs in groups.values():\n"
        "        evs.sort(key=lambda e: e['date'])\n\n"
        "        # Parte en sub-sesiones en cada START\n"
        "        sub_sessions, current = [], []\n"
        "        for ev in evs:\n"
        "            if ev['type_event'] == 'START' and current:\n"
        "                sub_sessions.append(current); current = []\n"
        "            current.append(ev)\n"
        "        if current: sub_sessions.append(current)\n\n"
        "        # Métricas por sub-sesión\n"
        "        for sub in sub_sessions:\n"
        "            total_buffer = sum(e['buffer_time'] for e in sub\n"
        "                               if e.get('buffer_time') is not None)\n"
        "            rebuffer_count = sum(1 for e in sub\n"
        "                if e['type_event'] == 'RE-BUFFERING')\n"
        "            completion_pct = round(\n"
        "                min(max_position / duration * 100, 100.0), 2)\n"
        "            type_events = {e['type_event'] for e in sub}\n"
        "            reached_complete = 'COMPLETE' in type_events\n\n"
        "# Escritura idempotente por la clave de negocio\n"
        "UpdateOne(\n"
        "    {'subscriber_id': s['subscriber_id'],\n"
        "     'customer_id':   s['customer_id'],\n"
        "     'start_time':    s['start_time']},\n"
        "    {'$set': s}, upsert=True)"
    )
    add_code_box(slide, code, Cm(16), Cm(4.5), Cm(16.5), Cm(13),
                 size=9, title="db/pipelines/sessions.py")

    add_footer(slide, page, total)


def slide_pipeline_content_stats(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Pipeline: content_stats",
                       "De sesiones a métricas agregadas por contenido (customer_id)",
                       number=12)

    add_card(slide, "Lo que calcula por contenido",
             [
                 "total_plays = nº de sesiones del contenido",
                 "unique_viewers = subscriber_ids únicos",
                 "rate(field) = % de sesiones que alcanzan el hito",
                 "  → firstquartile_rate, midpoint_rate, …, completion_rate",
                 "avg_buffer_time = promedio de total_buffer_time por sesión",
                 "rebuffer_rate = % de sesiones con rebuffer_count > 0",
                 "avg_completion_pct = progreso medio (capado a 100%)",
                 "plays_by_device = {ANDR: N, IOS: N, WEB: N}",
                 "last_updated = datetime.now(timezone.utc)",
             ],
             Cm(1.5), Cm(4.5), Cm(13.5), Cm(8.5),
             title_color=RED, body_size=11)

    add_card(slide, "Por qué pre-agregar",
             [
                 "El dashboard NUNCA recalcula en cada request",
                 "Un customer_id = una pieza puntual (un episodio)",
                 "upsert por customer_id mantiene el doc al día",
                 "Re-construye al ingerir nuevos CSVs (admin/upload-csv)",
             ],
             Cm(1.5), Cm(13.2), Cm(13.5), Cm(4.2),
             title_color=RED_DEEP, body_size=11)

    code = (
        "def _build_content_stats(sessions):\n"
        "    groups = defaultdict(list)\n"
        "    for s in sessions:\n"
        "        groups[s['customer_id']].append(s)\n\n"
        "    for customer_id, sess_list in groups.items():\n"
        "        n = len(sess_list)\n"
        "        first = sess_list[0]\n"
        "        unique_viewers = len(\n"
        "            {s['subscriber_id'] for s in sess_list})\n\n"
        "        def rate(field):  # % que alcanza un hito\n"
        "            return round(\n"
        "                sum(1 for s in sess_list if s.get(field))\n"
        "                / n * 100, 2)\n\n"
        "        # Buffer y rebuffer\n"
        "        buffers = [s['total_buffer_time'] for s in sess_list\n"
        "                   if s.get('total_buffer_time') is not None]\n"
        "        avg_buffer = round(sum(buffers) / len(buffers), 3)\n"
        "        rebuffer_sessions = sum(\n"
        "            1 for s in sess_list if s.get('rebuffer_count', 0) > 0)\n"
        "        rebuffer_rate = round(rebuffer_sessions / n * 100, 2)\n\n"
        "        docs.append({\n"
        "            'customer_id':      customer_id,\n"
        "            'total_plays':      n,\n"
        "            'unique_viewers':   unique_viewers,\n"
        "            'completion_rate':  rate('reached_complete'),\n"
        "            'avg_buffer_time':  avg_buffer,\n"
        "            'rebuffer_rate':    rebuffer_rate,\n"
        "            'plays_by_device':  device_counts,\n"
        "            'last_updated':     datetime.now(timezone.utc),\n"
        "        })"
    )
    add_code_box(slide, code, Cm(16), Cm(4.5), Cm(16.5), Cm(13),
                 size=9, title="db/pipelines/content_stats.py")

    add_footer(slide, page, total)


def slide_alertas(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Alertas en vivo  -  tres niveles",
                       "Detección de QoE degradada sobre la ventana de datos disponible",
                       number=13)

    levels = [
        ("ROJA  ·  Buffer crítico",
         RED,
         [
             "Sesión con total_buffer_time > 15 segundos",
             "En la ventana [earliest_event, latest_event]",
             "Ordena por buffer descendente (peor primero)",
             "Find proyectado: solo trae los campos necesarios",
         ]),
        ("AMARILLA  ·  Rebuffer rate",
         AMBER,
         [
             "Contenido con >30% de sesiones con rebuffer",
             "Mínimo 10 sesiones para descartar outliers",
             "Pipeline: $group por customer_id, $cond para contar",
             "$divide para la tasa, $match para filtrar",
         ]),
        ("AZUL  ·  Pause storm",
         INFO,
         [
             "Usuario con >5 PAUSEs en <8 minutos",
             "Find de PAUSEs ordenado por (subscriber, customer, date)",
             "Ventana deslizante en Python por usuario+contenido",
             "Una alerta por (subscriber, customer) — break tras detectar",
         ]),
    ]

    for i, (title, color, items) in enumerate(levels):
        x = Cm(1.5) + i * Cm(10.6)
        y = Cm(4.5)
        add_rect(slide, x, y, Cm(10), Cm(7), SURFACE, BORDER)
        add_rect(slide, x, y, Cm(10), Cm(0.85), color)
        add_text(slide, title, x, y + Cm(0.18), Cm(10), Cm(0.7),
                 size=14, bold=True, color=BLACK if color == AMBER else WHITE,
                 align=PP_ALIGN.CENTER)
        for j, it in enumerate(items):
            add_text(slide, "·  " + it, x + Cm(0.4),
                     y + Cm(1.3) + j * Cm(1.35), Cm(9.2), Cm(1.3),
                     size=11, color=TEXT)

    code = (
        "# Alerta amarilla — pipeline real\n"
        "pipeline = [\n"
        "    {\"$match\": {\"start_time\": {\"$gte\": start, \"$lte\": end}}},\n"
        "    {\"$group\": {\n"
        "        \"_id\":   \"$customer_id\",\n"
        "        \"title\": {\"$first\": \"$title\"},\n"
        "        \"total\": {\"$sum\": 1},\n"
        "        \"with_rebuffer\": {\n"
        "            \"$sum\": {\"$cond\":\n"
        "                [{\"$gt\": [\"$rebuffer_count\", 0]}, 1, 0]}}\n"
        "    }},\n"
        "    {\"$project\": {\"rate\": {\"$divide\": [\"$with_rebuffer\", \"$total\"]},\n"
        "                  \"title\": 1, \"total\": 1}},\n"
        "    {\"$match\": {\"rate\": {\"$gt\": 0.30}, \"total\": {\"$gte\": 10}}},\n"
        "    {\"$sort\":  {\"rate\": -1}},\n"
        "]"
    )
    add_code_box(slide, code, Cm(1.5), Cm(12.0), SLIDE_W - Cm(3), Cm(5.5),
                 size=10, title="db/repositories/alerts.py · get_rebuffer_rate_alerts()")

    add_footer(slide, page, total)


def slide_operadores(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Operadores de agregación",
                       "Todo el cálculo de las métricas vive dentro de MongoDB",
                       number=14)

    ops = [
        ("$group + $sum / $avg / $max",
         RED,
         "Agrega del lado del servidor sin traerse documentos completos",
         "overview (device-ranking, top-content) · qoe (startup_time) · users-map"),
        ("$switch",
         RED_DEEP,
         "Clasifica cada usuario por nº de eventos en una sola pasada",
         "users/user-profiles → Vague < 10 ≤ Mid < 50 ≤ Heavy"),
        ("$hour / $dayOfWeek",
         AMBER,
         "Extrae componentes de fecha para series temporales",
         "overview/activity-by-hour (24h) y users/activity-heatmap (24h × 7d)"),
        ("$cond + $divide",
         INFO,
         "Contadores condicionales y tasas calculadas in-pipeline",
         "alertas (rebuffer rate) · qoe ($round sobre promedios)"),
        ("$project + $match (post-group)",
         WHITE,
         "Reforma documentos y filtra DESPUÉS de agrupar",
         "content-completion-ranking filtra por total_plays mínimo"),
        ("$first / $addToSet",
         SUCCESS,
         "Trae el primer valor por grupo o un set único",
         "top-content (content_type del primer doc) · users-map (country_code)"),
    ]

    for i, (op, color, what, where) in enumerate(ops):
        y = Cm(4.5) + i * Cm(2.0)
        add_rect(slide, Cm(1.5), y, Cm(30.5), Cm(1.7), SURFACE, BORDER)
        add_rect(slide, Cm(1.5), y, Cm(0.2), Cm(1.7), color)
        add_text(slide, op, Cm(2.0), y + Cm(0.2), Cm(11), Cm(1.4),
                 size=15, bold=True, color=color, font_name="Consolas")
        add_text(slide, what, Cm(13.5), y + Cm(0.1), Cm(13), Cm(0.8),
                 size=11, color=TEXT)
        add_text(slide, where, Cm(13.5), y + Cm(0.85), Cm(13), Cm(0.8),
                 size=10, color=TEXT_DIM, italic=True)
        add_pill(slide, "EN USO", Cm(27), y + Cm(0.5), Cm(5), Cm(0.7),
                 color, BLACK if color in (WHITE, AMBER) else WHITE, 11)

    add_footer(slide, page, total)


def slide_change_streams(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Change Streams  -  los triggers",
                       "Dos tareas async escuchan a MongoDB y mantienen las derivadas al día",
                       number=15)

    code = (
        "async def watch_events_and_sync_sessions(db):\n"
        "    pipeline = [{\"$match\": {\"operationType\": \"insert\"}}]\n\n"
        "    async with db.events.watch(\n"
        "        pipeline, full_document=\"updateLookup\"\n"
        "    ) as stream:\n"
        "        async for change in stream:\n"
        "            doc = change[\"fullDocument\"]\n"
        "            subscriber = doc.get(\"subscriber_id\")\n"
        "            customer   = doc.get(\"customer_id\")\n\n"
        "            # Trae solo los eventos de ese usuario+contenido\n"
        "            events = await db.events.find(\n"
        "                {\"subscriber_id\": subscriber,\n"
        "                 \"customer_id\":   customer},\n"
        "                {proyección…}\n"
        "            ).sort(\"date\", 1).to_list(length=None)\n\n"
        "            # Reconstruye sesión(es) y upsert\n"
        "            for s in _build_sessions(events):\n"
        "                await db.sessions.update_one(\n"
        "                    _clave_unica(s),\n"
        "                    {\"$set\": s}, upsert=True)"
    )
    add_code_box(slide, code, Cm(1.5), Cm(4.5), Cm(19.5), Cm(13),
                 size=10, title="db/collections/sessions.py · watch_events_and_sync_sessions()")

    add_card(slide, "Arranque y resiliencia",
             [
                 "Lanzados en el lifespan() de FastAPI",
                 "asyncio.create_task() en background",
                 'pipeline filtra operationType="insert"',
                 'full_document="updateLookup" trae el doc completo',
             ],
             Cm(22.0), Cm(4.5), Cm(10.5), Cm(5.5),
             title_color=RED, body_size=11)

    add_card(slide, "Tarea hermana",
             [
                 "watch_sessions_and_sync_content_stats",
                 'opera con "insert" Y "update"',
                 "  (las sesiones se actualizan, no solo insertan)",
                 "Trae todas las sesiones del customer_id y reconstruye",
             ],
             Cm(22.0), Cm(10.3), Cm(10.5), Cm(3.5),
             title_color=RED_DEEP, body_size=11)

    add_card(slide, "Degradación con gracia",
             [
                 "_safe_watch() en api/main.py",
                 "Si el cluster no soporta change streams",
                 "  (standalone, no replica set) → loguea y sigue",
                 "Las derivadas se mantienen vía scripts/test_pipeline.py",
             ],
             Cm(22.0), Cm(14.1), Cm(10.5), Cm(3.4),
             title_color=WHITE, body_size=11)

    add_footer(slide, page, total)


def slide_api(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "API REST  -  18 endpoints, 5 routers",
                       "FastAPI con OpenAPI auto · response_model con Pydantic v2 · Depends(get_db)",
                       number=16)

    routers = [
        ("/api/overview",
         RED,
         [
             "GET /unique-users",
             "GET /total-plays",
             "GET /device-ranking",
             "GET /activity-by-hour",
             "GET /top-content?limit=N",
             "GET /users-map?limit=N",
         ]),
        ("/api/qoe",
         RED_DEEP,
         [
             "GET /buffer-by-content?limit=N",
             "GET /rebuffering-rate",
             "GET /startup-time",
             "GET /event-ranking",
         ]),
        ("/api/users",
         AMBER,
         [
             "GET /retention-funnel",
             "GET /activity-heatmap",
             "GET /content-completion-ranking",
             "GET /user-profiles",
         ]),
        ("/api/alerts",
         INFO,
         [
             "GET /  (rojo · amarillo · azul)",
             "Conteo por severidad",
             "Ventana de datos en period_start/end",
         ]),
        ("/api/admin",
         SUCCESS,
         [
             "GET /collections",
             "POST /upload-csv (multipart)",
             "GET /documents/{coll}?page=&page_size=",
             "PUT /documents/{coll}/{id}",
             "DELETE /documents/{coll}/{id}",
         ]),
    ]

    cols = 3
    rows = 2
    for i, (path, color, items) in enumerate(routers):
        col = i % cols
        row = i // cols
        x = Cm(1.5) + col * Cm(10.6)
        y = Cm(4.5) + row * Cm(6.3)
        add_rect(slide, x, y, Cm(10), Cm(5.8), SURFACE, BORDER)
        add_rect(slide, x, y, Cm(10), Cm(0.85), color)
        add_text(slide, path, x, y + Cm(0.18), Cm(10), Cm(0.7),
                 size=14, bold=True,
                 color=BLACK if color in (AMBER, WHITE) else WHITE,
                 align=PP_ALIGN.CENTER, font_name="Consolas")
        for j, it in enumerate(items):
            add_text(slide, "·  " + it, x + Cm(0.4),
                     y + Cm(1.15) + j * Cm(0.75), Cm(9.2), Cm(0.7),
                     size=11, color=TEXT, font_name="Consolas")

    # Bloque inferior con detalles transversales
    add_card(slide, "Transversales",
             [
                 "Health check en /api/health · OpenAPI en /docs",
                 "CORS abierto en dev, restringido a API_BASE_URL en prod",
                 "Lifespan: connect/disconnect + arranque de change streams",
                 "Todos los endpoints usan Depends(get_db) — no toca Mongo directo",
             ],
             Cm(23), Cm(11), Cm(9.5), Cm(6),
             title_color=RED, body_size=11)

    add_footer(slide, page, total)


def slide_dash_overview(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Dashboard · Overview",
                       "Plays, usuarios, mapa de burbujas, ranking de dispositivos y actividad por hora",
                       number=17)

    components = [
        ("KPI · Total plays",
         "28.584 sesiones de reproducción", RED),
        ("KPI · Usuarios únicos",
         "10.812 subscribers distintos", RED_DEEP),
        ("Top Content",
         "Ranking por total_plays con tipo de contenido", AMBER),
        ("Mapa de usuarios",
         "Leaflet · burbuja por usuario · centroide país + jitter md5", INFO),
        ("Device Ranking",
         "IOS · ANDR · WEB · ANDRTV · APPLETV · ROKU", SUCCESS),
        ("Activity by Hour",
         "Barras de 24 horas (gaps a 0 garantizados)", WHITE),
    ]

    cols = 3
    for i, (title, desc, color) in enumerate(components):
        col = i % cols
        row = i // cols
        x = Cm(1.5) + col * Cm(10.6)
        y = Cm(4.5) + row * Cm(5.5)
        add_rect(slide, x, y, Cm(10), Cm(5), SURFACE, BORDER)
        add_rect(slide, x, y, Cm(10), Cm(0.5), color)
        add_text(slide, title, x, y + Cm(0.8), Cm(10), Cm(1.0),
                 size=16, bold=True,
                 color=BLACK if color in (WHITE, AMBER) else color,
                 align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + Cm(0.5), y + Cm(2.2), Cm(9), Cm(2.5),
                 size=12, color=TEXT, align=PP_ALIGN.CENTER)

    add_rect(slide, Cm(1.5), Cm(15.7), SLIDE_W - Cm(3), Cm(1.8), SURFACE, RED, 1.0)
    add_text(slide,
             "Privacidad: el subscriber_id se publica SOLO con sus primeros 8 caracteres. "
             "La latitud/longitud usa el centroide del país + un jitter determinístico (md5 del id), "
             "así que un mismo usuario cae siempre en el mismo punto sin exponer ubicación real.",
             Cm(2.0), Cm(15.9), SLIDE_W - Cm(4), Cm(1.5),
             size=11, color=TEXT, italic=True)

    add_footer(slide, page, total)


def slide_dash_qoe(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Dashboard · QoE",
                       "Quality of Experience: buffer, rebuffering y startup time",
                       number=18)

    add_card(slide, "Buffer por contenido",
             [
                 "Top peores contenidos por avg_buffer_time",
                 "Pipeline sobre content_stats con $avg + $round",
                 "Filtra avg_buffer_time > 0 (descarta sin datos)",
                 "Devuelve title, avg_buffer_time, rebuffer_rate, total_plays",
             ],
             Cm(1.5), Cm(4.5), Cm(15.5), Cm(5.7),
             title_color=RED, body_size=11)

    add_card(slide, "Rebuffering rate",
             [
                 "Doble medida: por evento Y por sesión",
                 "Evento: # RE-BUFFERING / # total de eventos",
                 "Sesión: # sesiones con ≥1 rebuffer / # total sesiones",
                 "Real: ~10.36% eventos, ~8.84% sesiones",
             ],
             Cm(17.5), Cm(4.5), Cm(15), Cm(5.7),
             title_color=RED_DEEP, body_size=11)

    add_card(slide, "Startup time",
             [
                 "Mide buffer_time de los eventos START-BUFFERING",
                 "(buffering previo al primer frame de la sesión)",
                 "Promedio, máximo, conteo y distribución por buckets",
                 "Buckets: 0-1s · 1-2s · 2-5s · 5-10s · 10s+",
                 "Visualización tipo histograma horizontal (Recharts)",
             ],
             Cm(1.5), Cm(10.6), Cm(15.5), Cm(6.5),
             title_color=AMBER, body_size=11)

    add_card(slide, "Event ranking",
             [
                 "Conteo global de tipos de evento sobre events",
                 "Real: PAUSE 86.534 · RESUME 77.628",
                 "       RE-BUFFERING 31.089 · START 28.584",
                 "Permite ver el peso relativo de cada tipo en el log",
             ],
             Cm(17.5), Cm(10.6), Cm(15), Cm(6.5),
             title_color=INFO, body_size=11)

    add_footer(slide, page, total)


def slide_dash_users(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Dashboard · Usuarios y Alertas",
                       "Perfiles, funnel, heatmap, completion y alertas QoE",
                       number=19)

    add_card(slide, "Perfiles  (sobre events, $switch)",
             [
                 "Heavy:  ≥ 50 eventos   ·  ~835 usuarios",
                 "Mid:    10 ≤ x < 50    ·  ~2.146 usuarios",
                 "Vague:  < 10 eventos   ·  resto",
                 "Pipeline: $group por subscriber → $switch → $group",
                 "Devuelve count, total_events y avg_events por bucket",
             ],
             Cm(1.5), Cm(4.5), Cm(15.5), Cm(6.5),
             title_color=RED, body_size=11)

    add_card(slide, "Funnel de retención",
             [
                 "count_documents por hito sobre sessions:",
                 "  total · firstquartile · midpoint · thirdquartile · complete",
                 "Reales: 28.584 · 5.617 · 4.498 · 3.667 · 2.804",
                 "Tasa de finalización agregada: ~9.8%",
                 "Visualización en funnel Recharts (escalado de barras)",
             ],
             Cm(17.5), Cm(4.5), Cm(15), Cm(6.5),
             title_color=RED_DEEP, body_size=11)

    add_card(slide, "Heatmap 24h × 7d",
             [
                 "$group con {hour: $hour, weekday: $dayOfWeek}",
                 "168 celdas garantizadas (rellena con 0)",
                 "Mongo convention: 1=domingo … 7=sábado",
             ],
             Cm(1.5), Cm(11.4), Cm(10.5), Cm(6),
             title_color=AMBER, body_size=11)

    add_card(slide, "Completion ranking",
             [
                 "Top más completados y top más abandonados",
                 "Pipeline $group por title + $match total_plays ≥ N",
                 "Reduce ruido de contenidos con pocas vistas",
                 "Default min_plays = 5",
             ],
             Cm(12.5), Cm(11.4), Cm(10.0), Cm(6),
             title_color=INFO, body_size=11)

    add_card(slide, "Alertas (página dedicada)",
             [
                 "AlertsCounter: roja / amarilla / azul",
                 "AlertsTable: cada alerta con type, description,",
                 "  severity, timestamp, subscriber/customer",
             ],
             Cm(23), Cm(11.4), Cm(9.5), Cm(6),
             title_color=SUCCESS, body_size=11)

    add_footer(slide, page, total)


def slide_dash_admin(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Dashboard · Admin",
                       "Subida de CSVs y edición/borrado de documentos en vivo",
                       number=20)

    add_card(slide, "CollectionCounts",
             [
                 "GET /api/admin/collections",
                 "Conteo de events / sessions / content_stats",
                 "Total agregado de las tres",
                 "Útil para verificar la carga después de un upload",
             ],
             Cm(1.5), Cm(4.5), Cm(15.5), Cm(5.5),
             title_color=RED, body_size=11)

    add_card(slide, "CsvUploader",
             [
                 "POST /api/admin/upload-csv (multipart/form-data)",
                 "Validación: extensión .csv y body no vacío",
                 "Encadena: parse → insert → build_sessions → build_content_stats",
                 "Responde rows_total, rows_skipped, inserted, sessions_built, ...",
                 "El dashboard refleja los datos nuevos de inmediato",
             ],
             Cm(17.5), Cm(4.5), Cm(15), Cm(5.5),
             title_color=RED_DEEP, body_size=11)

    add_card(slide, "DocumentTable (CRUD en vivo)",
             [
                 "GET /api/admin/documents/{collection}?page=&page_size=",
                 "Paginación con DEFAULT_PAGE_SIZE=50, MAX_PAGE_SIZE=500",
                 "PUT /documents/{coll}/{id} aplica $set (_id es inmutable)",
                 "DELETE /documents/{coll}/{id} con ObjectId validado",
                 "_check_collection() rechaza cualquier colección fuera de ADMIN_COLLECTIONS",
                 "_stringify_id() convierte ObjectId → str para JSON",
             ],
             Cm(1.5), Cm(10.4), Cm(31), Cm(7),
             title_color=AMBER, body_size=12)

    add_footer(slide, page, total)


def slide_tests(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Tests, calidad y mitigaciones",
                       "104 tests en total: 81 backend (pytest) + 23 frontend (vitest)",
                       number=21)

    add_card(slide, "Backend · pytest + mongomock",
             [
                 "tests/db/test_pipelines.py  ·  reconstrucción de sesiones",
                 "tests/db/test_load_csv.py  ·  parse_row, cast, filtrado",
                 "tests/db/test_repositories.py  ·  aggregations de overview/qoe/users",
                 "tests/api/test_routes.py  ·  contrato de las 18 rutas",
                 "tests/api/test_admin.py  ·  upload CSV + CRUD",
                 "Mongo en memoria → no toca Atlas en tests",
                 "Resultado: 81 / 81 PASS",
             ],
             Cm(1.5), Cm(4.5), Cm(15.5), Cm(8),
             title_color=RED, body_size=11)

    add_card(slide, "Frontend · Vitest + Testing Library",
             [
                 "9 archivos de test · 23 tests",
                 "App.test.jsx — navegación entre páginas",
                 "Componentes presentacionales (KPICard, ContentRanking, …)",
                 "Mock global de fetch en beforeEach",
                 "jsdom como environment",
                 "Resultado: 23 / 23 PASS",
             ],
             Cm(17.5), Cm(4.5), Cm(15), Cm(8),
             title_color=RED_DEEP, body_size=11)

    add_card(slide, "Cosas que podrían salir mal (mitigaciones)",
             [
                 "Duplicados en sessions  →  índice ÚNICO + upsert por clave de negocio",
                 "Memoria en los find  →  pipelines proyectan campos · agregaciones en servidor",
                 "Bundle del front > 500 KB  →  warning conocido (recharts + leaflet)",
                 "Cluster sin change streams  →  _safe_watch() degrada y loguea",
                 "Unidades mixtas en Position/Duration  →  completion_pct capado a 100%",
             ],
             Cm(1.5), Cm(12.8), SLIDE_W - Cm(3), Cm(4.7),
             title_color=WHITE, body_size=12)

    add_footer(slide, page, total)


def slide_numeros(prs, total, page):
    slide = add_blank_slide(prs)
    add_section_header(slide, "Números del proyecto",
                       "Lo que hay corriendo en Atlas en este momento", number=None)

    nums = [
        ("300k",  "eventos (events)", RED),
        ("28.5k", "sesiones (sessions)", RED_DEEP),
        ("10.8k", "contenidos (content_stats)", AMBER),
        ("10.8k", "usuarios únicos", INFO),
        ("3",     "colecciones", WHITE),
        ("16",    "índices", SUCCESS),
        ("18",    "endpoints REST", RED),
        ("5",     "routers FastAPI", RED_DEEP),
        ("5",     "páginas del dashboard", AMBER),
        ("3",     "niveles de alerta", INFO),
        ("13",    "países en el mapa", SUCCESS),
        ("104",   "tests (81 back + 23 front)", RED),
    ]
    cols = 4
    for i, (num, label, color) in enumerate(nums):
        col = i % cols
        row = i // cols
        x = Cm(1.5) + col * Cm(8)
        y = Cm(4.5) + row * Cm(4.2)
        add_rect(slide, x, y, Cm(7.5), Cm(3.8), SURFACE, color, 1.5)
        add_text(slide, num, x, y + Cm(0.2), Cm(7.5), Cm(2.2),
                 size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + Cm(0.3), y + Cm(2.6), Cm(6.9), Cm(1.1),
                 size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER, italic=True)

    add_footer(slide, page, total)


def slide_closing(prs, total):
    slide = add_blank_slide(prs)
    add_rect(slide, Cm(0), Cm(0), Cm(2.2), SLIDE_H, RED)
    add_rect(slide, Cm(0), Cm(0.6), SLIDE_W, Cm(0.08), RED)

    add_text(slide, "WIN", Cm(3.3), Cm(3.0), Cm(10), Cm(4.5),
             size=110, bold=True, color=RED)
    add_text(slide, "SPORTS", Cm(12.5), Cm(3.0), Cm(20), Cm(4.5),
             size=110, bold=True, color=WHITE)

    add_text(slide,
             '"Que MongoDB haga el trabajo pesado: agregar dentro de la base, no fuera."',
             Cm(3.5), Cm(9.0), Cm(28), Cm(1.5),
             size=20, color=RED, italic=True)

    pts = [
        "Un log plano de eventos se vuelve sesiones y métricas por contenido",
        "Los pipelines de agregación reemplazan decenas de líneas de Python",
        "Change Streams mantienen las derivadas sincronizadas sin polling",
        "Arquitectura modular: API por repositorios, front solo consume REST",
        "104 tests respaldan el contrato; el bundle queda en un solo paso",
    ]
    for i, t in enumerate(pts):
        add_text(slide, "·  " + t, Cm(4.5), Cm(11.0) + i * Cm(0.9),
                 Cm(28), Cm(0.85), size=14, color=TEXT)

    add_text(slide, "Gracias", Cm(3.5), SLIDE_H - Cm(3.5), Cm(28), Cm(2),
             size=44, bold=True, color=RED, align=PP_ALIGN.CENTER)

    add_footer(slide, total, total)


# ─── Build ───────────────────────────────────────────────────────────────────


def build():
    prs = make_presentation()

    total = 26  # portada + disclaimer + agenda + 21 secciones + números + cierre

    slide_cover(prs, total)
    slide_legal_disclaimer(prs, total, 2)
    slide_agenda(prs, total, 3)
    slide_problema(prs, total, 4)
    slide_stack(prs, total, 5)
    slide_arquitectura(prs, total, 6)
    slide_modelo_overview(prs, total, 7)
    slide_coleccion_events(prs, total, 8)
    slide_coleccion_sessions(prs, total, 9)
    slide_coleccion_content_stats(prs, total, 10)
    slide_validators(prs, total, 11)
    slide_indices(prs, total, 12)
    slide_ingesta_csv(prs, total, 13)
    slide_pipeline_sessions(prs, total, 14)
    slide_pipeline_content_stats(prs, total, 15)
    slide_alertas(prs, total, 16)
    slide_operadores(prs, total, 17)
    slide_change_streams(prs, total, 18)
    slide_api(prs, total, 19)
    slide_dash_overview(prs, total, 20)
    slide_dash_qoe(prs, total, 21)
    slide_dash_users(prs, total, 22)
    slide_dash_admin(prs, total, 23)
    slide_tests(prs, total, 24)
    slide_numeros(prs, total, 25)
    slide_closing(prs, total)  # footer renderiza total/total = 26/26

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)

    # Si el archivo está abierto en PowerPoint, Windows lo bloquea para escritura.
    # En vez de fallar, escribimos a un nombre alterno con timestamp para que
    # la generación NUNCA quede sin entregar archivo justo antes de una demo.
    out = OUT_PATH
    try:
        prs.save(str(out))
    except PermissionError:
        from datetime import datetime
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        out = OUT_PATH.with_name(f"{OUT_PATH.stem}_{stamp}.pptx")
        prs.save(str(out))
        print("AVISO  -  el .pptx principal está abierto (probablemente en PowerPoint).")
        print(f"         Guardado en versión alterna: {out.name}")
    print(f"OK  -  presentación creada ({out.stat().st_size / 1024:.1f} KB)")
    print(f"OUT  -  {out}")


if __name__ == "__main__":
    build()
